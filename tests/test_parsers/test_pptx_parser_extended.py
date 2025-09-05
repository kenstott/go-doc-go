"""
Extended test suite for PPTX document parser to improve coverage.

Tests follow CLAUDE.md guidelines:
- Test design objectives, not implementation  
- Use pytest markers for categorization (unit, integration, performance)
- Follow DRY principles with reusable helpers
- Validate against ElementType and RelationshipType enums
"""

import json
import os
import tempfile
import time
from pathlib import Path
from typing import Dict, Any, List
from unittest.mock import Mock, patch, MagicMock

import pytest

# Import the parser and related modules
from go_doc_go.document_parser.pptx import PptxParser

try:
    import pptx
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Import enums for validation
from go_doc_go.storage.element_element import ElementType
from go_doc_go.relationships.structural import RelationshipType


# =============================================================================
# Test Helpers and Validators (DRY Principle)
# =============================================================================

def assert_valid_element(element: Dict[str, Any]) -> None:
    """Validate element structure against design specifications."""
    # Required fields
    required_fields = ["element_id", "doc_id", "element_type", "content_preview", "content_location"]
    for field in required_fields:
        assert field in element, f"Missing required field: {field}"
    
    # Validate element type against enum
    element_type = element["element_type"]
    valid_types = [e.value for e in ElementType]
    assert element_type in valid_types, f"Invalid element_type: {element_type}"


def assert_valid_relationship(relationship: Dict[str, Any]) -> None:
    """Validate relationship structure against design specifications."""
    required_fields = ["source_id", "target_id", "relationship_type", "relationship_id"]
    for field in required_fields:
        assert field in relationship, f"Missing required field: {field}"
    
    # Validate relationship type against enum
    rel_type = relationship["relationship_type"]
    valid_types = [r.value for r in RelationshipType]
    assert rel_type in valid_types, f"Invalid relationship_type: {rel_type}"


def assert_valid_parse_result(result: Dict[str, Any]) -> None:
    """Validate overall parse result structure."""
    assert "document" in result
    assert "elements" in result  
    assert "relationships" in result
    
    # Validate document structure
    doc = result["document"]
    assert "doc_id" in doc
    assert "doc_type" in doc
    assert "metadata" in doc
    
    # Validate all elements
    for element in result["elements"]:
        assert_valid_element(element)
    
    # Validate relationships
    for rel in result["relationships"]:
        assert_valid_relationship(rel)


@pytest.fixture
def temp_pptx_path():
    """Create a temporary PPTX file path."""
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
        temp_path = f.name
    yield temp_path
    try:
        os.unlink(temp_path)
    except FileNotFoundError:
        pass


@pytest.fixture
def simple_pptx_content(temp_pptx_path):
    """Create a simple PPTX for testing."""
    if not PPTX_AVAILABLE:
        pytest.skip("python-pptx not available")
        
    # Create a simple presentation
    prs = Presentation()
    
    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Test Presentation"
    subtitle.text = "Test Subtitle"
    
    # Add a content slide
    bullet_slide_layout = prs.slide_layouts[1]  
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Content Slide"
    tf = body_shape.text_frame
    tf.text = 'First bullet point'
    p = tf.add_paragraph()
    p.text = 'Second bullet point'
    
    prs.save(temp_pptx_path)
    
    return {
        "id": "/test.pptx",
        "binary_path": temp_pptx_path,  
        "metadata": {"doc_id": "pptx_test_123"}
    }


# =============================================================================
# Unit Tests for PPTX Parser
# =============================================================================

@pytest.mark.unit
class TestPptxParserExtended:
    """Extended tests to improve PPTX parser coverage."""

    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_missing_pptx_error_path(self):
        """Test the ImportError path when python-pptx is missing."""
        with patch('go_doc_go.document_parser.pptx.PPTX_AVAILABLE', False):
            with pytest.raises(ImportError, match="python-pptx is required"):
                PptxParser()

    def test_parser_initialization_config(self):
        """Test PPTX parser initialization with various configs."""
        # Skip if pptx not available
        if not PPTX_AVAILABLE:
            pytest.skip("python-pptx not available")
            
        # Default initialization
        parser1 = PptxParser()
        assert hasattr(parser1, 'config')
        
        # Custom configuration
        config = {
            "extract_images": True,
            "extract_notes": True, 
            "extract_charts": False,
            "max_content_preview": 200
        }
        parser2 = PptxParser(config)
        assert parser2.config == config

    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_parse_method_basic(self, simple_pptx_content):
        """Test basic parse method functionality."""
        parser = PptxParser()
        result = parser.parse(simple_pptx_content)
        
        assert_valid_parse_result(result)
        
        # Should have elements for slides and text
        elements = result["elements"]
        assert len(elements) > 0
        
        # Should have at least root and slide elements
        element_types = set(e["element_type"] for e in elements)
        assert ElementType.ROOT.value in element_types

    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available") 
    def test_extract_text_from_presentation(self, temp_pptx_path):
        """Test _extract_text_from_presentation method."""
        # Create test presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide Title"
        prs.save(temp_pptx_path)
        
        # Load and test text extraction
        prs_loaded = Presentation(temp_pptx_path)
        parser = PptxParser()
        
        text = parser._extract_text_from_presentation(prs_loaded)
        
        assert isinstance(text, str)
        assert "Test Slide Title" in text

    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_process_slide_method(self, temp_pptx_path):
        """Test _process_slide method."""
        # Create test slide with various content
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Slide Title"
        
        # Add text box
        text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
        text_frame = text_box.text_frame
        text_frame.text = "Text box content"
        
        prs.save(temp_pptx_path)
        
        # Load and test slide processing
        prs_loaded = Presentation(temp_pptx_path)
        slide_loaded = prs_loaded.slides[0]
        
        parser = PptxParser()
        elements, relationships = parser._process_slide(
            slide_loaded, "doc1", "root1", "source1", 0
        )
        
        assert isinstance(elements, list)
        assert isinstance(relationships, list)
        assert len(elements) > 0

    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_process_shape_method(self, temp_pptx_path):
        """Test _process_shape method."""
        # Create slide with shape
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add a simple rectangle shape
        from pptx.shapes.autoshape import Shape
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
        
        prs.save(temp_pptx_path)
        
        # Load and test shape processing
        prs_loaded = Presentation(temp_pptx_path)
        slide_loaded = prs_loaded.slides[0]
        shape_loaded = slide_loaded.shapes[0]
        
        parser = PptxParser()
        elements, relationships = parser._process_shape(
            shape_loaded, "doc1", "slide1", "source1", 0, 0
        )
        
        assert isinstance(elements, list)
        assert isinstance(relationships, list)

    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_extract_slide_notes(self, temp_pptx_path):
        """Test _extract_slide_notes method."""
        # Create slide with notes
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Slide with Notes"
        
        # Add notes
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = "These are speaker notes for the slide."
        
        prs.save(temp_pptx_path)
        
        # Load and test notes extraction
        prs_loaded = Presentation(temp_pptx_path)
        slide_loaded = prs_loaded.slides[0]
        
        parser = PptxParser()
        elements, relationships = parser._extract_slide_notes(
            slide_loaded, "doc1", "slide1", "source1"
        )
        
        assert isinstance(elements, list)
        assert isinstance(relationships, list)

    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_process_table_method(self, temp_pptx_path):
        """Test _process_table method."""
        # Create slide with table
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add a table
        rows, cols = 2, 3
        left = Inches(1)
        top = Inches(2)
        width = Inches(6)
        height = Inches(0.8)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Fill table with data
        table.cell(0, 0).text = 'Header 1'
        table.cell(0, 1).text = 'Header 2'
        table.cell(0, 2).text = 'Header 3'
        table.cell(1, 0).text = 'Data 1'
        table.cell(1, 1).text = 'Data 2'  
        table.cell(1, 2).text = 'Data 3'
        
        prs.save(temp_pptx_path)
        
        # Load and test table processing
        prs_loaded = Presentation(temp_pptx_path)
        slide_loaded = prs_loaded.slides[0]
        table_shape = slide_loaded.shapes[0]
        table_loaded = table_shape.table
        
        parser = PptxParser()
        elements, relationships = parser._process_table(
            table_loaded, "doc1", "slide1", "source1", 0
        )
        
        assert isinstance(elements, list)
        assert isinstance(relationships, list)
        assert len(elements) > 0

    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_extract_presentation_metadata(self, temp_pptx_path):
        """Test _extract_presentation_metadata static method."""
        # Create presentation with metadata
        prs = Presentation()
        
        # Set core properties
        prs.core_properties.title = "Test Presentation Title"
        prs.core_properties.author = "Test Author"
        prs.core_properties.comments = "Test Comments"
        
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Title Slide"
        
        prs.save(temp_pptx_path)
        
        # Load and test metadata extraction
        prs_loaded = Presentation(temp_pptx_path)
        base_metadata = {"existing": "data"}
        
        metadata = PptxParser._extract_presentation_metadata(prs_loaded, base_metadata)
        
        assert "existing" in metadata
        assert "title" in metadata or "author" in metadata

    def test_supports_location_method(self):
        """Test supports_location method."""
        if not PPTX_AVAILABLE:
            pytest.skip("python-pptx not available")
            
        parser = PptxParser()
        
        # Valid location for PPTX
        valid_location = {
            "slide": "1",
            "shape": "0"
        }
        assert parser.supports_location(valid_location) == True
        
        # Invalid location
        invalid_location = {
            "page": "1",
            "cell": "A1"
        }
        assert parser.supports_location(invalid_location) == False

    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_resolve_element_text(self, temp_pptx_path):
        """Test _resolve_element_text method."""
        # Create presentation with text
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Slide Title Text"
        prs.save(temp_pptx_path)
        
        parser = PptxParser()
        location_data = {
            "slide": "0",
            "shape": "0",
            "source": temp_pptx_path
        }
        
        with open(temp_pptx_path, 'rb') as f:
            binary_content = f.read()
            
        text = parser._resolve_element_text(location_data, binary_content)
        
        assert isinstance(text, str)

    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available") 
    def test_resolve_element_content(self, temp_pptx_path):
        """Test _resolve_element_content method."""
        # Create presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Content to Resolve"
        prs.save(temp_pptx_path)
        
        parser = PptxParser()
        location_data = {
            "slide": "0",
            "shape": "0",
            "source": temp_pptx_path
        }
        
        with open(temp_pptx_path, 'rb') as f:
            binary_content = f.read()
            
        content = parser._resolve_element_content(location_data, binary_content)
        
        assert isinstance(content, dict)
        assert "text" in content

    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_create_root_element(self):
        """Test _create_root_element method."""
        parser = PptxParser()
        
        root = parser._create_root_element("test_doc", "test_source")
        
        assert isinstance(root, dict)
        assert root["element_type"] == ElementType.ROOT.value
        assert root["doc_id"] == "test_doc"

    def test_generate_id_method(self):
        """Test _generate_id method."""
        if not PPTX_AVAILABLE:
            pytest.skip("python-pptx not available")
            
        parser = PptxParser()
        
        id1 = parser._generate_id("test_")
        id2 = parser._generate_id("test_")
        
        assert id1.startswith("test_")
        assert id2.startswith("test_")
        assert id1 != id2  # Should be unique

    def test_generate_hash_method(self):
        """Test _generate_hash method."""
        if not PPTX_AVAILABLE:
            pytest.skip("python-pptx not available")
            
        parser = PptxParser()
        
        hash1 = parser._generate_hash("test content")
        hash2 = parser._generate_hash("test content")
        hash3 = parser._generate_hash("different content")
        
        assert hash1 == hash2  # Same content = same hash
        assert hash1 != hash3  # Different content = different hash


# =============================================================================
# Integration Tests
# =============================================================================

@pytest.mark.integration
class TestPptxParserIntegration:
    """Integration tests for PPTX parser with real presentations."""

    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_complex_presentation_parsing(self, temp_pptx_path):
        """Test parsing a complex presentation with multiple content types."""
        # Create complex presentation
        prs = Presentation()
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "Complex Presentation"
        title_slide.shapes.placeholders[1].text = "With Multiple Content Types"
        
        # Content slide with bullets
        bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
        bullet_slide.shapes.title.text = "Bullet Points"
        tf = bullet_slide.shapes.placeholders[1].text_frame
        tf.text = "First bullet"
        p = tf.add_paragraph()
        p.text = "Second bullet"
        
        # Blank slide with custom content
        blank_slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add text box
        text_box = blank_slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Inches(1)
        )
        text_box.text_frame.text = "Custom text box content"
        
        prs.save(temp_pptx_path)
        
        # Parse the presentation
        content = {
            "id": "/complex.pptx",
            "binary_path": temp_pptx_path,
            "metadata": {"doc_id": "complex_test"}
        }
        
        parser = PptxParser()
        result = parser.parse(content)
        
        assert_valid_parse_result(result)
        
        # Should have multiple slides
        slide_elements = [e for e in result["elements"] if "slide" in e["element_type"]]
        assert len(slide_elements) >= 3  # At least 3 slides


# =============================================================================  
# Error Handling Tests
# =============================================================================

@pytest.mark.unit
class TestPptxParserErrorHandling:
    """Test error handling in PPTX parser."""

    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_missing_binary_path(self):
        """Test handling of missing binary_path."""
        parser = PptxParser()
        
        content = {
            "id": "/test.pptx",
            "metadata": {}
            # Missing binary_path
        }
        
        with pytest.raises((ValueError, KeyError)):
            parser.parse(content)

    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_invalid_file_path(self):
        """Test handling of invalid file path."""
        parser = PptxParser()
        
        content = {
            "id": "/test.pptx", 
            "binary_path": "/nonexistent/file.pptx",
            "metadata": {}
        }
        
        with pytest.raises(FileNotFoundError):
            parser.parse(content)

    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_corrupted_presentation_file(self, temp_pptx_path):
        """Test handling of corrupted PPTX file."""
        # Create a corrupted file (just text content)
        with open(temp_pptx_path, 'w') as f:
            f.write("This is not a valid PPTX file")
        
        content = {
            "id": "/corrupted.pptx",
            "binary_path": temp_pptx_path,
            "metadata": {}
        }
        
        parser = PptxParser()
        
        # Should raise an exception for corrupted file
        with pytest.raises(Exception):  # Could be various exceptions depending on implementation
            parser.parse(content)


if __name__ == "__main__":
    pytest.main([__file__, "-v"])